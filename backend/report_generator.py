"""
Professional YouTube Analytics Report Generator v3.0
YouTube 알고리즘 기반 컨설팅급 분석 보고서 생성기
PDF 내보내기 지원
"""
from datetime import datetime
import json
import io
import os

# PDF 생성 라이브러리 (선택적 - pdfkit 또는 weasyprint)
PDF_LIBRARY = None
PDFKIT_CONFIG = None

def find_wkhtmltopdf():
    """wkhtmltopdf 실행 파일 경로 자동 탐색"""
    import shutil

    # 1. PATH에서 찾기
    path_result = shutil.which('wkhtmltopdf')
    if path_result:
        return path_result

    # 2. 일반적인 Windows 설치 경로들
    common_paths = [
        r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe',
        r'C:\Program Files (x86)\wkhtmltopdf\bin\wkhtmltopdf.exe',
        r'C:\wkhtmltopdf\bin\wkhtmltopdf.exe',
        os.path.expanduser(r'~\wkhtmltopdf\bin\wkhtmltopdf.exe'),
        # Chocolatey 설치 경로
        r'C:\ProgramData\chocolatey\bin\wkhtmltopdf.exe',
        r'C:\tools\wkhtmltopdf\bin\wkhtmltopdf.exe',
    ]

    for path in common_paths:
        if os.path.exists(path):
            return path

    return None

try:
    import pdfkit
    WKHTMLTOPDF_PATH = find_wkhtmltopdf()
    if WKHTMLTOPDF_PATH:
        PDFKIT_CONFIG = pdfkit.configuration(wkhtmltopdf=WKHTMLTOPDF_PATH)
        PDF_LIBRARY = 'pdfkit'
        print(f"[INFO] pdfkit 사용 가능 (wkhtmltopdf: {WKHTMLTOPDF_PATH})")
    else:
        print("[WARN] wkhtmltopdf를 찾을 수 없습니다.")
        print("       설치 방법: choco install wkhtmltopdf")
        print("       또는: https://wkhtmltopdf.org/downloads.html")
except ImportError:
    pass

if PDF_LIBRARY is None:
    try:
        from weasyprint import HTML, CSS
        PDF_LIBRARY = 'weasyprint'
        print("[INFO] weasyprint 사용 가능")
    except (ImportError, OSError) as e:
        PDF_LIBRARY = None
        print(f"[WARN] PDF 라이브러리 없음 - PDF 생성 비활성화")
        print("       설치: pip install pdfkit && choco install wkhtmltopdf")


class ReportGenerator:
    """전문 분석 보고서 생성기 v3.0 - YouTube 알고리즘 기반"""

    THEME = {
        'primary': '#FF0000',
        'secondary': '#282828',
        'background': '#f9f9f9',
        'text': '#0f0f0f',
        'text_light': '#606060',
        'border': '#e5e5e5',
        'success': '#2e7d32',
        'warning': '#f57c00',
        'danger': '#c62828',
        'viral': '#FF0000',
        'hit': '#ff9800',
        'average': '#2196f3',
        'underperform': '#9e9e9e',
        'card_bg': '#ffffff',
    }

    def __init__(self, analysis_data: dict, channel_data: dict):
        self.analysis = analysis_data
        self.channel = channel_data

    def generate_html_report(self) -> str:
        """YouTube 알고리즘 기반 HTML 보고서 생성"""
        channel_summary = self.analysis.get('channel_summary', {})
        classification_stats = self.analysis.get('classification_stats', {})
        success = self.analysis.get('success_analysis', {})
        failure = self.analysis.get('failure_analysis', {})
        recommendations = self.analysis.get('recommendations', [])
        upload_patterns = self.analysis.get('upload_patterns', {})
        growth = self.analysis.get('growth_trends', {})
        trend_analysis = self.analysis.get('trend_analysis', {})
        algorithm_insights = self.analysis.get('algorithm_insights', {})
        detailed_metrics = self.analysis.get('detailed_metrics', {})
        benchmarks = self.analysis.get('youtube_benchmarks', {})

        html = f'''<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>YouTube 채널 분석 보고서 - {channel_summary.get('channel_name', '채널')}</title>
    <link href="https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Noto Sans KR', -apple-system, BlinkMacSystemFont, sans-serif;
            background: {self.THEME['background']};
            color: {self.THEME['text']};
            line-height: 1.7;
            font-size: 15px;
        }}
        .report-container {{ max-width: 1200px; margin: 0 auto; padding: 40px 20px; }}

        /* 표지 */
        .cover {{
            background: linear-gradient(135deg, {self.THEME['secondary']} 0%, #1a1a1a 100%);
            color: white;
            padding: 60px 40px;
            border-radius: 16px;
            margin-bottom: 40px;
            position: relative;
            overflow: hidden;
        }}
        .cover::before {{
            content: '';
            position: absolute;
            top: 0; right: 0;
            width: 300px; height: 300px;
            background: {self.THEME['primary']};
            opacity: 0.1;
            border-radius: 50%;
            transform: translate(50%, -50%);
        }}
        .cover-badge {{
            display: inline-block;
            background: {self.THEME['primary']};
            padding: 8px 16px;
            border-radius: 4px;
            font-size: 13px;
            font-weight: 600;
            margin-bottom: 20px;
        }}
        .cover h1 {{ font-size: 36px; font-weight: 700; margin-bottom: 12px; }}
        .cover-channel {{ font-size: 24px; opacity: 0.9; margin-bottom: 30px; }}
        .cover-meta {{ display: flex; gap: 30px; font-size: 14px; opacity: 0.7; flex-wrap: wrap; }}

        /* 섹션 */
        .section {{
            background: {self.THEME['card_bg']};
            border-radius: 12px;
            padding: 32px;
            margin-bottom: 24px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.06);
        }}
        .section-title {{
            font-size: 22px;
            font-weight: 700;
            color: {self.THEME['secondary']};
            margin-bottom: 24px;
            padding-bottom: 12px;
            border-bottom: 3px solid {self.THEME['primary']};
            display: flex;
            align-items: center;
            gap: 12px;
        }}

        /* 알고리즘 헬스 카드 */
        .algorithm-health {{
            background: linear-gradient(135deg, #FF0000 0%, #CC0000 100%);
            color: white;
            padding: 30px;
            border-radius: 12px;
            margin-bottom: 24px;
            text-align: center;
        }}
        .algorithm-health h3 {{ font-size: 18px; margin-bottom: 10px; opacity: 0.9; }}
        .algorithm-health .score {{ font-size: 48px; font-weight: 700; }}
        .algorithm-health .status {{ font-size: 20px; margin-top: 10px; }}

        /* 지표 그리드 */
        .metrics-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(180px, 1fr)); gap: 16px; }}
        .metric-card {{
            text-align: center;
            padding: 24px 16px;
            background: white;
            border-radius: 12px;
            box-shadow: 0 2px 8px rgba(0,0,0,0.04);
            border: 1px solid {self.THEME['border']};
        }}
        .metric-card .value {{
            font-size: 28px;
            font-weight: 700;
            color: {self.THEME['primary']};
            margin-bottom: 4px;
        }}
        .metric-card .label {{ font-size: 13px; color: {self.THEME['text_light']}; }}
        .metric-card .benchmark {{
            font-size: 11px;
            margin-top: 8px;
            padding: 4px 8px;
            border-radius: 4px;
            background: {self.THEME['background']};
        }}
        .metric-card .benchmark.good {{ background: #e8f5e9; color: #2e7d32; }}
        .metric-card .benchmark.warn {{ background: #fff3e0; color: #f57c00; }}
        .metric-card .benchmark.bad {{ background: #ffebee; color: #c62828; }}

        /* 분류 그리드 */
        .classification-grid {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 16px; margin-bottom: 30px; }}
        .class-card {{
            padding: 24px;
            border-radius: 12px;
            text-align: center;
            position: relative;
            overflow: hidden;
        }}
        .class-card.viral {{ background: linear-gradient(135deg, #ffebee, #fff); border: 2px solid {self.THEME['viral']}; }}
        .class-card.hit {{ background: linear-gradient(135deg, #fff3e0, #fff); border: 2px solid {self.THEME['hit']}; }}
        .class-card.average {{ background: linear-gradient(135deg, #e3f2fd, #fff); border: 2px solid {self.THEME['average']}; }}
        .class-card.underperform {{ background: linear-gradient(135deg, #fafafa, #fff); border: 2px solid {self.THEME['underperform']}; }}
        .class-card .emoji {{ font-size: 32px; margin-bottom: 8px; }}
        .class-card .count {{ font-size: 36px; font-weight: 700; }}
        .class-card .title {{ font-size: 14px; font-weight: 600; margin-top: 4px; }}
        .class-card .stats {{ font-size: 12px; color: {self.THEME['text_light']}; margin-top: 12px; padding-top: 12px; border-top: 1px solid rgba(0,0,0,0.1); }}

        /* 트렌드 박스 */
        .trend-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(250px, 1fr)); gap: 16px; }}
        .trend-box {{
            padding: 20px;
            background: {self.THEME['background']};
            border-radius: 12px;
            border-left: 4px solid;
        }}
        .trend-box.up {{ border-color: {self.THEME['success']}; }}
        .trend-box.down {{ border-color: {self.THEME['danger']}; }}
        .trend-box.stable {{ border-color: {self.THEME['warning']}; }}
        .trend-box h4 {{ font-size: 14px; color: {self.THEME['text_light']}; margin-bottom: 8px; }}
        .trend-box .values {{ display: flex; align-items: center; gap: 12px; }}
        .trend-box .old {{ color: {self.THEME['text_light']}; }}
        .trend-box .arrow {{ font-size: 20px; }}
        .trend-box .new {{ font-weight: 700; font-size: 18px; }}
        .trend-box .change {{ font-size: 14px; margin-top: 8px; }}
        .trend-box .change.positive {{ color: {self.THEME['success']}; }}
        .trend-box .change.negative {{ color: {self.THEME['danger']}; }}

        /* 비디오 카드 */
        .video-card {{
            display: grid;
            grid-template-columns: 120px 1fr;
            gap: 20px;
            padding: 20px;
            background: {self.THEME['background']};
            border-radius: 12px;
            margin-bottom: 16px;
            align-items: start;
        }}
        .video-thumb {{
            width: 120px;
            height: 68px;
            border-radius: 8px;
            background: #ddd;
            overflow: hidden;
        }}
        .video-thumb img {{ width: 100%; height: 100%; object-fit: cover; }}
        .video-info h4 {{ font-size: 15px; font-weight: 600; margin-bottom: 8px; line-height: 1.4; }}
        .video-metrics {{ display: flex; flex-wrap: wrap; gap: 16px; margin: 12px 0; }}
        .video-metric {{
            display: flex;
            flex-direction: column;
            font-size: 13px;
        }}
        .video-metric .val {{ font-weight: 700; color: {self.THEME['secondary']}; }}
        .video-metric .lbl {{ color: {self.THEME['text_light']}; font-size: 11px; }}
        .video-score {{
            display: inline-flex;
            align-items: center;
            gap: 8px;
            padding: 8px 16px;
            border-radius: 20px;
            font-size: 14px;
            font-weight: 600;
        }}
        .video-score.viral {{ background: {self.THEME['viral']}; color: white; }}
        .video-score.hit {{ background: {self.THEME['hit']}; color: white; }}
        .video-score.average {{ background: {self.THEME['average']}; color: white; }}
        .video-score.underperform {{ background: {self.THEME['underperform']}; color: white; }}
        .video-reasons {{
            margin-top: 12px;
            padding: 12px;
            background: white;
            border-radius: 8px;
            font-size: 13px;
        }}
        .video-reasons ul {{ list-style: none; padding: 0; margin: 0; }}
        .video-reasons li {{ padding: 4px 0; padding-left: 20px; position: relative; }}
        .video-reasons li::before {{ content: '→'; position: absolute; left: 0; color: {self.THEME['primary']}; }}

        /* 추천사항 */
        .recommendation {{
            padding: 24px;
            border-radius: 12px;
            margin-bottom: 16px;
            border-left: 5px solid;
        }}
        .recommendation.critical {{ background: linear-gradient(90deg, #ffebee, #fff); border-color: {self.THEME['danger']}; }}
        .recommendation.high {{ background: linear-gradient(90deg, #fff3e0, #fff); border-color: {self.THEME['warning']}; }}
        .recommendation.medium {{ background: linear-gradient(90deg, #e3f2fd, #fff); border-color: {self.THEME['average']}; }}
        .recommendation h4 {{ font-size: 16px; font-weight: 600; margin-bottom: 12px; }}
        .recommendation .targets {{ font-size: 13px; color: {self.THEME['text_light']}; margin-bottom: 12px; }}
        .recommendation ul {{ margin-left: 20px; }}
        .recommendation li {{ margin-bottom: 8px; }}

        /* 차트 */
        .chart-container {{ height: 280px; margin: 20px 0; }}
        .chart-row {{ display: grid; grid-template-columns: 1fr 1fr; gap: 24px; }}

        /* 분포 바 */
        .dist-bar {{ display: flex; height: 40px; border-radius: 8px; overflow: hidden; margin: 16px 0; }}
        .dist-segment {{ display: flex; align-items: center; justify-content: center; color: white; font-size: 12px; font-weight: 600; }}

        /* 테이블 */
        .data-table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
        .data-table th, .data-table td {{ padding: 14px; border: 1px solid {self.THEME['border']}; }}
        .data-table th {{ background: {self.THEME['secondary']}; color: white; font-weight: 600; }}
        .data-table tr:nth-child(even) {{ background: #f9f9f9; }}
        .data-table .score-cell {{ font-weight: 700; }}

        /* PDF 버튼 */
        .pdf-controls {{
            position: fixed;
            top: 20px;
            right: 20px;
            z-index: 1000;
        }}
        .pdf-btn {{
            background: {self.THEME['primary']};
            color: white;
            border: none;
            padding: 12px 24px;
            border-radius: 8px;
            font-size: 14px;
            font-weight: 600;
            cursor: pointer;
            box-shadow: 0 4px 12px rgba(255,0,0,0.3);
        }}
        .pdf-btn:hover {{ background: #cc0000; }}

        /* 페이지 구분 */
        .page-break {{ page-break-before: always; }}

        @media print {{
            body {{ background: white; }}
            .section {{ box-shadow: none; border: 1px solid {self.THEME['border']}; }}
            .pdf-controls {{ display: none !important; }}
            .cover {{ break-after: page; }}
        }}

        @media (max-width: 768px) {{
            .classification-grid {{ grid-template-columns: repeat(2, 1fr); }}
            .chart-row {{ grid-template-columns: 1fr; }}
            .video-card {{ grid-template-columns: 1fr; }}
            .video-thumb {{ width: 100%; height: auto; aspect-ratio: 16/9; }}
        }}
    </style>
</head>
<body>
<div class="report-container">

    <div class="pdf-controls">
        <button class="pdf-btn" onclick="window.print()">PDF 다운로드</button>
    </div>

    <!-- 표지 -->
    <div class="cover">
        <span class="cover-badge">YouTube Analytics Report v3.0</span>
        <h1>채널 심층 분석 보고서</h1>
        <p class="cover-channel">{channel_summary.get('channel_name', '채널명')}</p>
        <div class="cover-meta">
            <span>분석일: {datetime.now().strftime('%Y년 %m월 %d일')}</span>
            <span>분석 영상: {channel_summary.get('total_videos_analyzed', 0)}개</span>
            <span>구독자: {self._format_number(channel_summary.get('subscriber_count', 0))}</span>
        </div>
    </div>

    <!-- 알고리즘 헬스 체크 -->
    {self._generate_algorithm_health_section(algorithm_insights, channel_summary)}

    <!-- 핵심 지표 요약 -->
    <div class="section">
        <h2 class="section-title">핵심 지표 요약</h2>
        {self._generate_metrics_summary(channel_summary, benchmarks)}
    </div>

    <!-- 영상 성과 분류 -->
    <div class="section">
        <h2 class="section-title">영상 성과 분류</h2>
        {self._generate_classification_section(classification_stats)}
    </div>

    <!-- 트렌드 분석 -->
    <div class="section">
        <h2 class="section-title">트렌드 분석 (최근 vs 과거)</h2>
        {self._generate_trend_section(trend_analysis)}
    </div>

    <div class="page-break"></div>

    <!-- 성공 영상 분석 -->
    <div class="section">
        <h2 class="section-title">잘 터진 영상 분석</h2>
        {self._generate_success_section(success)}
    </div>

    <!-- 저조 영상 분석 -->
    <div class="section">
        <h2 class="section-title">저조한 영상 분석</h2>
        {self._generate_failure_section(failure)}
    </div>

    <div class="page-break"></div>

    <!-- 개선 추천사항 -->
    <div class="section">
        <h2 class="section-title">데이터 기반 개선 추천</h2>
        {self._generate_recommendations_section(recommendations)}
    </div>

    <!-- 상세 분포 -->
    <div class="section">
        <h2 class="section-title">상세 지표 분포</h2>
        {self._generate_distribution_section(detailed_metrics)}
    </div>

    <!-- 업로드 패턴 -->
    <div class="section">
        <h2 class="section-title">업로드 패턴 분석</h2>
        {self._generate_upload_section(upload_patterns)}
    </div>

    <!-- 전체 영상 목록 -->
    <div class="section">
        <h2 class="section-title">전체 영상 성과 상세</h2>
        {self._generate_all_videos_table(classification_stats)}
    </div>

    <!-- 푸터 -->
    <div style="text-align:center;padding:40px;color:{self.THEME['text_light']};font-size:13px;">
        <hr style="border:none;border-top:1px dashed {self.THEME['border']};margin-bottom:20px;">
        <p>본 보고서는 YouTube Data API 데이터와 2025-2026 알고리즘 분석 기준을 적용하여 생성되었습니다.</p>
        <p>생성일시: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</p>
    </div>

</div>

<script>
// 분류 차트
const classCtx = document.getElementById('classificationChart');
if (classCtx) {{
    new Chart(classCtx, {{
        type: 'doughnut',
        data: {{
            labels: ['바이럴', '히트', '평균', '저조'],
            datasets: [{{
                data: [{classification_stats.get('viral', {}).get('count', 0)}, {classification_stats.get('hit', {}).get('count', 0)}, {classification_stats.get('average', {}).get('count', 0)}, {classification_stats.get('underperform', {}).get('count', 0)}],
                backgroundColor: ['{self.THEME['viral']}', '{self.THEME['hit']}', '{self.THEME['average']}', '{self.THEME['underperform']}'],
                borderWidth: 0
            }}]
        }},
        options: {{
            responsive: true,
            maintainAspectRatio: false,
            plugins: {{
                legend: {{ position: 'right', labels: {{ font: {{ size: 14 }}, padding: 16 }} }}
            }},
            cutout: '60%'
        }}
    }});
}}
</script>

</body>
</html>'''
        return html

    def _generate_algorithm_health_section(self, insights: dict, summary: dict) -> str:
        """알고리즘 헬스 섹션"""
        health = insights.get('overall_health', '분석 중')
        signals = insights.get('algorithm_signals', [])

        # 평균 알고리즘 점수 계산
        avg_engagement = summary.get('avg_engagement_rate', 0)
        avg_like = summary.get('avg_like_ratio', 0)

        signal_html = ''
        for sig in signals[:4]:
            status_color = self.THEME['success'] if sig.get('status') == 'positive' else self.THEME['danger']
            signal_html += f'''
            <div style="display:flex;align-items:center;gap:8px;padding:8px 0;">
                <span style="color:{status_color};font-size:20px;">{'✓' if sig.get('status') == 'positive' else '✗'}</span>
                <span>{sig.get('message', '')}</span>
            </div>'''

        return f'''
        <div class="algorithm-health">
            <h3>YouTube 알고리즘 건강도</h3>
            <div class="status">{health}</div>
            <div style="margin-top:20px;display:flex;justify-content:center;gap:40px;">
                <div>
                    <div style="font-size:24px;font-weight:700;">{avg_engagement:.2f}%</div>
                    <div style="opacity:0.8;font-size:13px;">평균 참여율</div>
                </div>
                <div>
                    <div style="font-size:24px;font-weight:700;">{avg_like:.2f}%</div>
                    <div style="opacity:0.8;font-size:13px;">좋아요 비율</div>
                </div>
                <div>
                    <div style="font-size:24px;font-weight:700;">{summary.get('avg_view_velocity', 0):.0f}</div>
                    <div style="opacity:0.8;font-size:13px;">일평균 조회속도</div>
                </div>
            </div>
        </div>
        <div class="section" style="margin-top:-10px;">
            <h3 style="font-size:16px;margin-bottom:16px;">알고리즘 신호 분석</h3>
            {signal_html if signal_html else '<p style="color:#666;">신호 분석 데이터가 없습니다.</p>'}
        </div>'''

    def _generate_metrics_summary(self, summary: dict, benchmarks: dict) -> str:
        """핵심 지표 요약"""
        eng_rate = summary.get('avg_engagement_rate', 0)
        like_ratio = summary.get('avg_like_ratio', 0)

        # 벤치마크 상태
        eng_bench = benchmarks.get('engagement_rate', {})
        eng_status = 'good' if eng_rate >= eng_bench.get('good', 5) else 'warn' if eng_rate >= eng_bench.get('average', 3) else 'bad'

        like_bench = benchmarks.get('like_ratio', {})
        like_status = 'good' if like_ratio >= like_bench.get('good', 3) else 'warn' if like_ratio >= like_bench.get('average', 2) else 'bad'

        return f'''
        <div class="metrics-grid">
            <div class="metric-card">
                <div class="value">{self._format_number(summary.get('subscriber_count', 0))}</div>
                <div class="label">총 구독자</div>
            </div>
            <div class="metric-card">
                <div class="value">{summary.get('total_videos_analyzed', 0)}</div>
                <div class="label">분석 영상</div>
            </div>
            <div class="metric-card">
                <div class="value">{self._format_number(summary.get('avg_views_per_video', 0))}</div>
                <div class="label">평균 조회수</div>
            </div>
            <div class="metric-card">
                <div class="value">{summary.get('avg_view_velocity', 0):.0f}</div>
                <div class="label">조회 속도 (일평균)</div>
            </div>
            <div class="metric-card">
                <div class="value">{eng_rate:.2f}%</div>
                <div class="label">참여율</div>
                <div class="benchmark {eng_status}">{summary.get('engagement_benchmark_status', '')}</div>
            </div>
            <div class="metric-card">
                <div class="value">{like_ratio:.2f}%</div>
                <div class="label">좋아요 비율</div>
                <div class="benchmark {like_status}">{summary.get('like_ratio_benchmark_status', '')}</div>
            </div>
        </div>'''

    def _generate_classification_section(self, stats: dict) -> str:
        """성과 분류 섹션"""
        html = '<div class="classification-grid">'

        categories = [
            ('viral', '바이럴', '알고리즘 추천 영상'),
            ('hit', '히트', '좋은 성과 영상'),
            ('average', '평균', '채널 평균 수준'),
            ('underperform', '저조', '개선 필요 영상'),
        ]

        for cls, title, desc in categories:
            data = stats.get(cls, {})
            count = data.get('count', 0)
            avg_views = data.get('avg_views', 0)
            avg_velocity = data.get('avg_velocity', 0)
            avg_score = data.get('avg_score', 0)

            html += f'''
            <div class="class-card {cls}">
                <div class="count" style="color:{self.THEME[cls]}">{count}</div>
                <div class="title">{title}</div>
                <div class="stats">
                    조회 {self._format_number(avg_views)}<br>
                    속도 {avg_velocity:.0f}/일<br>
                    점수 {avg_score:.0f}
                </div>
            </div>'''

        html += '</div>'

        # 차트
        html += '''
        <div class="chart-container" style="height:250px;">
            <canvas id="classificationChart"></canvas>
        </div>'''

        return html

    def _generate_trend_section(self, trends: dict) -> str:
        """트렌드 분석 섹션"""
        if trends.get('message'):
            return f'<p style="color:{self.THEME["text_light"]};padding:20px;text-align:center;">{trends["message"]}</p>'

        html = '<div class="trend-grid">'

        # 조회 속도 트렌드
        vv = trends.get('view_velocity', {})
        if vv:
            trend_class = 'up' if vv.get('trend') == '상승' else 'down' if vv.get('trend') == '하락' else 'stable'
            change_class = 'positive' if vv.get('change_percent', 0) > 0 else 'negative'
            html += f'''
            <div class="trend-box {trend_class}">
                <h4>조회 속도 (핵심!)</h4>
                <div class="values">
                    <span class="old">{vv.get('older_avg', 0):.0f}/일</span>
                    <span class="arrow">→</span>
                    <span class="new">{vv.get('recent_avg', 0):.0f}/일</span>
                </div>
                <div class="change {change_class}">{vv.get('change_percent', 0):+.1f}% {vv.get('interpretation', '')}</div>
            </div>'''

        # 참여율 트렌드
        eng = trends.get('engagement', {})
        if eng:
            trend_class = 'up' if eng.get('trend') == '상승' else 'down' if eng.get('trend') == '하락' else 'stable'
            change_class = 'positive' if eng.get('change_percent', 0) > 0 else 'negative'
            html += f'''
            <div class="trend-box {trend_class}">
                <h4>참여율</h4>
                <div class="values">
                    <span class="old">{eng.get('older_avg', 0):.2f}%</span>
                    <span class="arrow">→</span>
                    <span class="new">{eng.get('recent_avg', 0):.2f}%</span>
                </div>
                <div class="change {change_class}">{eng.get('change_percent', 0):+.1f}% ({eng.get('benchmark_status', '')})</div>
            </div>'''

        # 제목 길이 트렌드
        title_len = trends.get('title_length', {})
        if title_len:
            trend_class = 'stable'
            html += f'''
            <div class="trend-box {trend_class}">
                <h4>제목 길이</h4>
                <div class="values">
                    <span class="old">{title_len.get('older_avg', 0):.0f}자</span>
                    <span class="arrow">→</span>
                    <span class="new">{title_len.get('recent_avg', 0):.0f}자</span>
                </div>
                <div class="change">{title_len.get('optimal_range', '')}</div>
            </div>'''

        # 제목 CTR 점수 트렌드
        ctr = trends.get('title_ctr_score', {})
        if ctr:
            trend_class = 'up' if ctr.get('trend') == '개선' else 'down' if ctr.get('trend') == '악화' else 'stable'
            html += f'''
            <div class="trend-box {trend_class}">
                <h4>제목 CTR 점수</h4>
                <div class="values">
                    <span class="old">{ctr.get('older_avg', 0):.0f}점</span>
                    <span class="arrow">→</span>
                    <span class="new">{ctr.get('recent_avg', 0):.0f}점</span>
                </div>
                <div class="change">{ctr.get('change', 0):+.0f}점 변화</div>
            </div>'''

        # 성공률 트렌드
        success = trends.get('success_rate', {})
        if success:
            trend_class = 'up' if success.get('trend') == '개선' else 'down' if success.get('trend') == '악화' else 'stable'
            change_class = 'positive' if success.get('change', 0) > 0 else 'negative'
            html += f'''
            <div class="trend-box {trend_class}">
                <h4>성공 영상 비율</h4>
                <div class="values">
                    <span class="old">{success.get('older', 0):.0f}%</span>
                    <span class="arrow">→</span>
                    <span class="new">{success.get('recent', 0):.0f}%</span>
                </div>
                <div class="change {change_class}">{success.get('change', 0):+.1f}%p</div>
            </div>'''

        html += '</div>'
        return html

    def _generate_success_section(self, success: dict) -> str:
        """성공 영상 섹션"""
        if success.get('message'):
            return f'<p style="color:{self.THEME["text_light"]};padding:20px;">{success["message"]}</p>'

        html = f'''
        <div class="metrics-grid" style="margin-bottom:30px;">
            <div class="metric-card">
                <div class="value">{success.get('total_count', 0)}</div>
                <div class="label">성공 영상 수</div>
            </div>
            <div class="metric-card">
                <div class="value">{self._format_number(success.get('avg_views', 0))}</div>
                <div class="label">평균 조회수</div>
            </div>
            <div class="metric-card">
                <div class="value">{success.get('avg_velocity', 0):.0f}</div>
                <div class="label">평균 조회속도</div>
            </div>
            <div class="metric-card">
                <div class="value">{success.get('avg_engagement', 0):.2f}%</div>
                <div class="label">평균 참여율</div>
            </div>
            <div class="metric-card">
                <div class="value">{success.get('avg_score', 0):.0f}</div>
                <div class="label">평균 알고리즘 점수</div>
            </div>
        </div>

        <h3 style="font-size:16px;margin-bottom:16px;">TOP 성과 영상</h3>'''

        for video in success.get('top_videos', [])[:5]:
            cls = video.get('classification', 'hit')
            reasons = video.get('success_reasons', [])
            breakdown = video.get('score_breakdown', {})
            title_info = video.get('title_analysis', {})

            html += f'''
            <div class="video-card">
                <div class="video-thumb">
                    <img src="{video.get('thumbnail_url', '')}" alt="thumbnail" onerror="this.style.display='none'">
                </div>
                <div class="video-info">
                    <h4>{video.get('title', '')}</h4>
                    <div class="video-metrics">
                        <div class="video-metric">
                            <span class="val">{self._format_number(video.get('view_count', 0))}</span>
                            <span class="lbl">조회수</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{video.get('view_velocity', 0):.0f}/일</span>
                            <span class="lbl">조회속도</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{video.get('engagement_rate', 0):.2f}%</span>
                            <span class="lbl">참여율</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{video.get('like_ratio', 0):.2f}%</span>
                            <span class="lbl">좋아요율</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{title_info.get('score', 0)}점</span>
                            <span class="lbl">제목CTR</span>
                        </div>
                    </div>
                    <span class="video-score {cls}">알고리즘 점수: {video.get('algorithm_score', 0):.0f}</span>
                    <div class="video-reasons">
                        <strong>성공 요인:</strong>
                        <ul>{''.join(f"<li>{r}</li>" for r in reasons)}</ul>
                    </div>
                </div>
            </div>'''

        # 성공 패턴
        patterns = success.get('success_patterns', [])
        if patterns:
            html += f'''
            <div style="margin-top:24px;padding:20px;background:{self.THEME['background']};border-radius:12px;border-left:4px solid {self.THEME['success']};">
                <h4 style="color:{self.THEME['success']};margin-bottom:12px;">성공 영상 공통 패턴</h4>
                <ul style="margin-left:20px;">{''.join(f"<li style='margin-bottom:6px;'>{p}</li>" for p in patterns)}</ul>
            </div>'''

        return html

    def _generate_failure_section(self, failure: dict) -> str:
        """저조 영상 섹션"""
        if failure.get('message'):
            return f'<p style="color:{self.THEME["text_light"]};padding:20px;">{failure["message"]}</p>'

        html = f'''
        <div class="metrics-grid" style="margin-bottom:30px;">
            <div class="metric-card">
                <div class="value">{failure.get('total_count', 0)}</div>
                <div class="label">저조 영상 수</div>
            </div>
            <div class="metric-card">
                <div class="value">{self._format_number(failure.get('avg_views', 0))}</div>
                <div class="label">평균 조회수</div>
            </div>
            <div class="metric-card">
                <div class="value">{failure.get('avg_velocity', 0):.0f}</div>
                <div class="label">평균 조회속도</div>
            </div>
            <div class="metric-card">
                <div class="value">{failure.get('avg_engagement', 0):.2f}%</div>
                <div class="label">평균 참여율</div>
            </div>
        </div>

        <h3 style="font-size:16px;margin-bottom:16px;">저조 영상 상세</h3>'''

        for video in failure.get('bottom_videos', [])[:5]:
            reasons = video.get('failure_reasons', [])
            title_info = video.get('title_analysis', {})

            html += f'''
            <div class="video-card">
                <div class="video-thumb">
                    <img src="{video.get('thumbnail_url', '')}" alt="thumbnail" onerror="this.style.display='none'">
                </div>
                <div class="video-info">
                    <h4>{video.get('title', '')}</h4>
                    <div class="video-metrics">
                        <div class="video-metric">
                            <span class="val">{self._format_number(video.get('view_count', 0))}</span>
                            <span class="lbl">조회수</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{video.get('view_velocity', 0):.0f}/일</span>
                            <span class="lbl">조회속도</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{video.get('engagement_rate', 0):.2f}%</span>
                            <span class="lbl">참여율</span>
                        </div>
                        <div class="video-metric">
                            <span class="val">{title_info.get('score', 0)}점</span>
                            <span class="lbl">제목CTR</span>
                        </div>
                    </div>
                    <span class="video-score underperform">알고리즘 점수: {video.get('algorithm_score', 0):.0f}</span>
                    <div class="video-reasons">
                        <strong>부진 원인:</strong>
                        <ul>{''.join(f"<li>{r}</li>" for r in reasons)}</ul>
                    </div>
                </div>
            </div>'''

        # 실패 패턴
        patterns = failure.get('failure_patterns', [])
        if patterns:
            html += f'''
            <div style="margin-top:24px;padding:20px;background:#ffebee;border-radius:12px;border-left:4px solid {self.THEME['danger']};">
                <h4 style="color:{self.THEME['danger']};margin-bottom:12px;">개선 필요 사항</h4>
                <ul style="margin-left:20px;">{''.join(f"<li style='margin-bottom:6px;'>{p}</li>" for p in patterns)}</ul>
            </div>'''

        # 성공 vs 실패 비교
        comparison = failure.get('comparison_with_success', {})
        if comparison:
            html += f'''
            <div style="margin-top:24px;">
                <h4 style="margin-bottom:16px;">성공 영상 vs 저조 영상 비교</h4>
                <table class="data-table">
                    <thead>
                        <tr>
                            <th>지표</th>
                            <th>성공 영상</th>
                            <th>저조 영상</th>
                            <th>차이</th>
                        </tr>
                    </thead>
                    <tbody>'''

            metrics = [
                ('view_velocity', '조회 속도', '/일'),
                ('engagement_rate', '참여율', '%'),
                ('title_length', '제목 길이', '자'),
                ('title_ctr_score', '제목 CTR 점수', '점'),
            ]

            for key, label, unit in metrics:
                data = comparison.get(key, {})
                if data:
                    success_val = data.get('success_avg', 0)
                    fail_val = data.get('failure_avg', 0)
                    diff = data.get('difference', success_val - fail_val)
                    diff_color = self.THEME['success'] if diff > 0 else self.THEME['danger']

                    html += f'''
                        <tr>
                            <td>{label}</td>
                            <td style="text-align:right;">{success_val:.1f}{unit}</td>
                            <td style="text-align:right;">{fail_val:.1f}{unit}</td>
                            <td style="text-align:right;color:{diff_color};font-weight:600;">{diff:+.1f}{unit}</td>
                        </tr>'''

            html += '</tbody></table></div>'

        return html

    def _generate_recommendations_section(self, recommendations: list) -> str:
        """추천사항 섹션"""
        if not recommendations:
            return '<p style="color:#666;padding:20px;">추천사항을 생성할 데이터가 부족합니다.</p>'

        html = ''
        for rec in recommendations:
            priority = rec.get('priority', 'medium')
            category = rec.get('category', '')
            current = rec.get('current', '')
            target = rec.get('target', '')
            suggestions = rec.get('suggestions', [])

            targets_html = ''
            if current and target:
                targets_html = f'<div class="targets">현재: {current} → 목표: {target}</div>'

            html += f'''
            <div class="recommendation {priority}">
                <h4>{category}</h4>
                {targets_html}
                <ul>{''.join(f"<li>{s}</li>" for s in suggestions)}</ul>
            </div>'''

        return html

    def _generate_distribution_section(self, metrics: dict) -> str:
        """분포 섹션"""
        html = '<div style="display:grid;grid-template-columns:repeat(auto-fit, minmax(300px, 1fr));gap:24px;">'

        # 알고리즘 점수 분포
        score_dist = metrics.get('performance_score_distribution', {})
        if score_dist:
            html += '<div>'
            html += '<h4 style="margin-bottom:12px;">알고리즘 점수 분포</h4>'
            total = sum(score_dist.values()) or 1
            html += '<div class="dist-bar">'
            colors = [self.THEME['danger'], self.THEME['warning'], self.THEME['average'], self.THEME['hit'], self.THEME['viral']]
            for i, (label, count) in enumerate(score_dist.items()):
                pct = count / total * 100
                if pct > 0:
                    html += f'<div class="dist-segment" style="flex:{pct};background:{colors[i % len(colors)]}">{count}</div>'
            html += '</div>'
            for label, count in score_dist.items():
                html += f'<div style="font-size:13px;padding:4px 0;">{label}: {count}개</div>'
            html += '</div>'

        # 참여율 분포
        eng_dist = metrics.get('engagement_distribution', {})
        if eng_dist:
            html += '<div>'
            html += '<h4 style="margin-bottom:12px;">참여율 분포</h4>'
            total = sum(eng_dist.values()) or 1
            html += '<div class="dist-bar">'
            colors = [self.THEME['underperform'], self.THEME['average'], self.THEME['hit'], self.THEME['success'], self.THEME['viral']]
            for i, (label, count) in enumerate(eng_dist.items()):
                pct = count / total * 100
                if pct > 0:
                    html += f'<div class="dist-segment" style="flex:{pct};background:{colors[i % len(colors)]}">{count}</div>'
            html += '</div>'
            for label, count in eng_dist.items():
                html += f'<div style="font-size:13px;padding:4px 0;">{label}: {count}개</div>'
            html += '</div>'

        # 조회 속도 분포
        vv_dist = metrics.get('view_velocity_distribution', {})
        if vv_dist:
            html += '<div>'
            html += '<h4 style="margin-bottom:12px;">조회 속도 분포</h4>'
            total = sum(vv_dist.values()) or 1
            html += '<div class="dist-bar">'
            colors = [self.THEME['underperform'], self.THEME['average'], self.THEME['hit'], self.THEME['warning'], self.THEME['viral']]
            for i, (label, count) in enumerate(vv_dist.items()):
                pct = count / total * 100
                if pct > 0:
                    html += f'<div class="dist-segment" style="flex:{pct};background:{colors[i % len(colors)]}">{count}</div>'
            html += '</div>'
            for label, count in vv_dist.items():
                html += f'<div style="font-size:13px;padding:4px 0;">{label}: {count}개</div>'
            html += '</div>'

        html += '</div>'
        return html

    def _generate_upload_section(self, patterns: dict) -> str:
        """업로드 패턴 섹션"""
        if patterns.get('message'):
            return f'<p style="color:{self.THEME["text_light"]};padding:20px;">{patterns["message"]}</p>'

        weekday_dist = patterns.get('weekday_distribution', {})
        day_korean = {
            'Monday': '월', 'Tuesday': '화', 'Wednesday': '수',
            'Thursday': '목', 'Friday': '금', 'Saturday': '토', 'Sunday': '일'
        }

        html = f'''
        <div class="metrics-grid" style="margin-bottom:24px;">
            <div class="metric-card">
                <div class="value">{patterns.get('avg_upload_interval_days', 0):.1f}일</div>
                <div class="label">평균 업로드 간격</div>
            </div>
            <div class="metric-card">
                <div class="value">{patterns.get('upload_frequency', '분석 중')}</div>
                <div class="label">업로드 빈도</div>
            </div>
        </div>

        <h4 style="margin-bottom:12px;">요일별 업로드 분포</h4>
        <div style="display:flex;gap:12px;flex-wrap:wrap;">'''

        for day, count in weekday_dist.items():
            html += f'''
            <div style="flex:1;min-width:80px;text-align:center;padding:16px;background:{self.THEME['background']};border-radius:8px;">
                <div style="font-size:24px;font-weight:700;color:{self.THEME['primary']}">{count}</div>
                <div style="font-size:13px;color:{self.THEME['text_light']}">{day_korean.get(day, day)}</div>
            </div>'''

        html += '</div>'
        return html

    def _generate_all_videos_table(self, stats: dict) -> str:
        """전체 영상 테이블"""
        all_videos = []
        for cls in ['viral', 'hit', 'average', 'underperform']:
            all_videos.extend(stats.get(cls, {}).get('videos', []))

        all_videos = sorted(all_videos, key=lambda x: x.get('algorithm_score', 0), reverse=True)

        if not all_videos:
            return '<p style="color:#666;">표시할 영상이 없습니다.</p>'

        class_labels = {
            'viral': ('바이럴', self.THEME['viral']),
            'hit': ('히트', self.THEME['hit']),
            'average': ('평균', self.THEME['average']),
            'underperform': ('저조', self.THEME['underperform'])
        }

        html = '''
        <div style="overflow-x:auto;">
        <table class="data-table">
            <thead>
                <tr>
                    <th>#</th>
                    <th>제목</th>
                    <th style="text-align:right;">조회수</th>
                    <th style="text-align:right;">속도</th>
                    <th style="text-align:right;">참여율</th>
                    <th style="text-align:right;">좋아요율</th>
                    <th style="text-align:center;">분류</th>
                    <th style="text-align:center;">점수</th>
                </tr>
            </thead>
            <tbody>'''

        for i, video in enumerate(all_videos[:30], 1):
            cls = video.get('classification', 'average')
            label, color = class_labels.get(cls, ('평균', self.THEME['average']))

            # 분류에 따라 video 딕셔너리에서 classification 키를 찾음
            for c_name, c_data in stats.items():
                for v in c_data.get('videos', []):
                    if v.get('video_id') == video.get('video_id'):
                        cls = c_name
                        label, color = class_labels.get(cls, ('평균', self.THEME['average']))
                        break

            html += f'''
                <tr>
                    <td>{i}</td>
                    <td>
                        <a href="https://youtube.com/watch?v={video.get('video_id', '')}" target="_blank" style="color:{self.THEME['text']};text-decoration:none;">
                            {video.get('title', '')[:45]}{'...' if len(video.get('title', '')) > 45 else ''}
                        </a>
                    </td>
                    <td style="text-align:right;">{self._format_number(video.get('view_count', 0))}</td>
                    <td style="text-align:right;">{video.get('view_velocity', 0):.0f}/일</td>
                    <td style="text-align:right;">{video.get('engagement_rate', 0):.2f}%</td>
                    <td style="text-align:right;">{video.get('like_ratio', 0):.2f}%</td>
                    <td style="text-align:center;">
                        <span style="background:{color};color:white;padding:4px 10px;border-radius:10px;font-size:12px;">{label}</span>
                    </td>
                    <td style="text-align:center;" class="score-cell">{video.get('algorithm_score', 0):.0f}</td>
                </tr>'''

        html += '</tbody></table></div>'
        return html

    def _format_number(self, num: int) -> str:
        """숫자 포맷팅"""
        if num >= 1000000:
            return f'{num / 1000000:.1f}M'
        elif num >= 1000:
            return f'{num / 1000:.1f}K'
        return str(int(num))

    def generate_summary_text(self) -> str:
        """텍스트 요약"""
        summary = self.analysis.get('channel_summary', {})
        growth = self.analysis.get('growth_trends', {})
        insights = self.analysis.get('algorithm_insights', {})

        return f"""
YouTube 채널 분석 요약 (v3.0)
{'='*50}
채널명: {summary.get('channel_name', '')}
분석일: {datetime.now().strftime('%Y-%m-%d')}

알고리즘 건강도: {insights.get('overall_health', '분석 중')}

핵심 지표:
- 구독자: {self._format_number(summary.get('subscriber_count', 0))}
- 분석 영상: {summary.get('total_videos_analyzed', 0)}개
- 평균 조회수: {self._format_number(summary.get('avg_views_per_video', 0))}
- 조회 속도: {summary.get('avg_view_velocity', 0):.0f}/일
- 참여율: {summary.get('avg_engagement_rate', 0):.2f}% ({summary.get('engagement_benchmark_status', '')})
- 좋아요 비율: {summary.get('avg_like_ratio', 0):.2f}% ({summary.get('like_ratio_benchmark_status', '')})
- 성장률: {growth.get('growth_rate_percent', 0):+.1f}%
"""

    def generate_pdf_report(self, output_path: str = None) -> bytes:
        """PDF 보고서 생성 (pdfkit 또는 weasyprint 사용)

        Args:
            output_path: 파일 저장 경로 (None이면 bytes 반환)

        Returns:
            PDF bytes (output_path가 None일 때)
        """
        if PDF_LIBRARY is None:
            raise ImportError("PDF 라이브러리가 필요합니다. pip install pdfkit (+ wkhtmltopdf 설치)")

        # HTML 생성 (PDF용으로 스타일 조정)
        html_content = self._generate_pdf_optimized_html()

        if PDF_LIBRARY == 'pdfkit':
            # pdfkit 사용 (wkhtmltopdf 필요)
            # 한글 폰트 지원을 위한 설정
            options = {
                'page-size': 'A4',
                'margin-top': '10mm',
                'margin-right': '10mm',
                'margin-bottom': '10mm',
                'margin-left': '10mm',
                'encoding': 'UTF-8',
                'enable-local-file-access': '',
                'print-media-type': '',
                'no-outline': '',
                'dpi': 300,
                'image-quality': 100,
                'disable-smart-shrinking': '',
            }

            if output_path:
                pdfkit.from_string(html_content, output_path, options=options, configuration=PDFKIT_CONFIG)
                return None
            else:
                pdf_bytes = pdfkit.from_string(html_content, False, options=options, configuration=PDFKIT_CONFIG)
                return pdf_bytes

        elif PDF_LIBRARY == 'weasyprint':
            # weasyprint 사용
            from weasyprint import HTML, CSS

            pdf_css = CSS(string='''
                @page { size: A4; margin: 10mm; }
                body {
                    font-family: "Malgun Gothic", "맑은 고딕", sans-serif !important;
                    font-size: 10pt;
                    line-height: 1.4;
                    color: #000 !important;
                }
            ''')

            html_doc = HTML(string=html_content)

            if output_path:
                html_doc.write_pdf(output_path, stylesheets=[pdf_css])
                return None
            else:
                pdf_buffer = io.BytesIO()
                html_doc.write_pdf(pdf_buffer, stylesheets=[pdf_css])
                return pdf_buffer.getvalue()

    def generate_pptx_report(self) -> bytes:
        """PPT 보고서 생성 (한글 완벽 지원)"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE

        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        channel_summary = self.analysis.get('channel_summary', {})
        classification_stats = self.analysis.get('classification_stats', {})
        success = self.analysis.get('success_analysis', {})
        failure = self.analysis.get('failure_analysis', {})
        recommendations = self.analysis.get('recommendations', [])
        algorithm_insights = self.analysis.get('algorithm_insights', {})
        trend_analysis = self.analysis.get('trend_analysis', {})

        def add_title_slide(title, subtitle):
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # 배경색
            background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(33, 33, 33)
            background.line.fill.background()

            # 빨간색 배지
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(3), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(244, 67, 54)
            badge.line.fill.background()
            tf = badge.text_frame
            tf.paragraphs[0].text = "YouTube Analytics Report v3.0"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # 타이틀
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12), Inches(1))
            tf = title_box.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(44)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # 서브타이틀
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.6))
            tf = sub_box.text_frame
            tf.paragraphs[0].text = subtitle
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.color.rgb = RGBColor(189, 189, 189)

            # 날짜
            date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.4))
            tf = date_box.text_frame
            tf.paragraphs[0].text = f"분석일: {datetime.now().strftime('%Y년 %m월 %d일')}"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(158, 158, 158)

            return slide

        def add_content_slide(title, content_func):
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)

            # 타이틀 영역
            title_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(33, 33, 33)
            title_bg.line.fill.background()

            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
            tf = title_box.text_frame
            tf.paragraphs[0].text = title
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            content_func(slide)
            return slide

        def add_metric_box(slide, x, y, w, h, value, label, bg_color):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            box.fill.solid()
            box.fill.fore_color.rgb = bg_color
            box.line.fill.background()

            # 값
            val_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.3), Inches(w), Inches(0.6))
            tf = val_box.text_frame
            tf.paragraphs[0].text = str(value)
            tf.paragraphs[0].font.size = Pt(32)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # 라벨
            lbl_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.9), Inches(w), Inches(0.4))
            tf = lbl_box.text_frame
            tf.paragraphs[0].text = label
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.color.rgb = RGBColor(97, 97, 97)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 슬라이드 1: 표지
        add_title_slide(
            channel_summary.get('channel_name', '채널 분석 보고서'),
            f"구독자 {self._format_number(channel_summary.get('subscriber_count', 0))} | 분석 영상 {channel_summary.get('total_videos_analyzed', 0)}개"
        )

        # 슬라이드 2: 핵심 지표
        def metrics_content(slide):
            metrics = [
                (self._format_number(channel_summary.get('subscriber_count', 0)), "구독자", RGBColor(255, 243, 224)),
                (self._format_number(channel_summary.get('avg_views_per_video', 0)), "평균 조회수", RGBColor(232, 245, 233)),
                (f"{channel_summary.get('avg_view_velocity', 0):.0f}/일", "조회 속도", RGBColor(227, 242, 253)),
                (f"{channel_summary.get('avg_engagement_rate', 0):.2f}%", "참여율", RGBColor(243, 229, 245)),
            ]
            for i, (val, label, color) in enumerate(metrics):
                add_metric_box(slide, 0.5 + i * 3.1, 1.8, 2.9, 1.5, val, label, color)

            # 알고리즘 건강도
            health = algorithm_insights.get('overall_health', '분석 중')
            health_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3.8), Inches(12.3), Inches(1.2))
            health_box.fill.solid()
            health_box.fill.fore_color.rgb = RGBColor(103, 58, 183)
            health_box.line.fill.background()

            hb = slide.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.8))
            tf = hb.text_frame
            tf.paragraphs[0].text = f"YouTube 알고리즘 건강도: {health}"
            tf.paragraphs[0].font.size = Pt(28)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_content_slide("핵심 지표 요약", metrics_content)

        # 슬라이드 3: 영상 분류
        def classification_content(slide):
            class_data = [
                ('viral', '바이럴', RGBColor(211, 47, 47), '알고리즘 추천'),
                ('hit', '히트', RGBColor(245, 124, 0), '좋은 성과'),
                ('average', '평균', RGBColor(25, 118, 210), '채널 평균'),
                ('underperform', '저조', RGBColor(117, 117, 117), '개선 필요'),
            ]
            for i, (cls, label, color, desc) in enumerate(class_data):
                data = classification_stats.get(cls, {})
                count = data.get('count', 0)
                avg_views = data.get('avg_views', 0)
                avg_score = data.get('avg_score', 0)

                x = 0.5 + i * 3.2
                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.6), Inches(3), Inches(2.5))
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(250, 250, 250)
                box.line.color.rgb = color

                # 개수
                cnt_box = slide.shapes.add_textbox(Inches(x), Inches(1.8), Inches(3), Inches(0.8))
                tf = cnt_box.text_frame
                tf.paragraphs[0].text = str(count)
                tf.paragraphs[0].font.size = Pt(48)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = color
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                # 라벨
                lbl_box = slide.shapes.add_textbox(Inches(x), Inches(2.6), Inches(3), Inches(0.4))
                tf = lbl_box.text_frame
                tf.paragraphs[0].text = label
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                # 상세
                detail_box = slide.shapes.add_textbox(Inches(x), Inches(3.1), Inches(3), Inches(0.8))
                tf = detail_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = f"평균 조회: {self._format_number(avg_views)}\n점수: {avg_score:.0f}"
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.color.rgb = RGBColor(97, 97, 97)
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_content_slide("영상 성과 분류", classification_content)

        # 슬라이드 4: TOP 성공 영상
        def success_content(slide):
            top_videos = success.get('top_videos', [])[:5]
            if not top_videos:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
                tf = tb.text_frame
                tf.paragraphs[0].text = "성공 영상 데이터가 없습니다."
                tf.paragraphs[0].font.size = Pt(18)
                return

            y = 1.5
            for i, v in enumerate(top_videos):
                title = v.get('title', '')[:50] + ('...' if len(v.get('title', '')) > 50 else '')
                views = self._format_number(v.get('view_count', 0))
                score = v.get('algorithm_score', 0)

                # 배경
                row_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1))
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = RGBColor(232, 245, 233)
                row_bg.line.fill.background()

                # 순위
                rank_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.25), Inches(0.5), Inches(0.5))
                tf = rank_box.text_frame
                tf.paragraphs[0].text = f"#{i+1}"
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)

                # 제목
                title_box = slide.shapes.add_textbox(Inches(1.4), Inches(y + 0.3), Inches(8), Inches(0.5))
                tf = title_box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

                # 조회수
                views_box = slide.shapes.add_textbox(Inches(9.5), Inches(y + 0.3), Inches(1.5), Inches(0.5))
                tf = views_box.text_frame
                tf.paragraphs[0].text = views
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

                # 점수
                score_box = slide.shapes.add_textbox(Inches(11.2), Inches(y + 0.3), Inches(1.3), Inches(0.5))
                tf = score_box.text_frame
                tf.paragraphs[0].text = f"{score:.0f}점"
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

                y += 1.1

        add_content_slide("✅ TOP 성공 영상", success_content)

        # 슬라이드 5: 저조 영상
        def failure_content(slide):
            bottom_videos = failure.get('bottom_videos', [])[:5]
            if not bottom_videos:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
                tf = tb.text_frame
                tf.paragraphs[0].text = "저조 영상 데이터가 없습니다."
                tf.paragraphs[0].font.size = Pt(18)
                return

            y = 1.5
            for i, v in enumerate(bottom_videos):
                title = v.get('title', '')[:50] + ('...' if len(v.get('title', '')) > 50 else '')
                views = self._format_number(v.get('view_count', 0))
                score = v.get('algorithm_score', 0)

                row_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1))
                row_bg.fill.solid()
                row_bg.fill.fore_color.rgb = RGBColor(255, 235, 238)
                row_bg.line.fill.background()

                rank_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.25), Inches(0.5), Inches(0.5))
                tf = rank_box.text_frame
                tf.paragraphs[0].text = f"#{i+1}"
                tf.paragraphs[0].font.size = Pt(20)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(198, 40, 40)

                title_box = slide.shapes.add_textbox(Inches(1.4), Inches(y + 0.3), Inches(8), Inches(0.5))
                tf = title_box.text_frame
                tf.paragraphs[0].text = title
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

                views_box = slide.shapes.add_textbox(Inches(9.5), Inches(y + 0.3), Inches(1.5), Inches(0.5))
                tf = views_box.text_frame
                tf.paragraphs[0].text = views
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

                score_box = slide.shapes.add_textbox(Inches(11.2), Inches(y + 0.3), Inches(1.3), Inches(0.5))
                tf = score_box.text_frame
                tf.paragraphs[0].text = f"{score:.0f}점"
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = RGBColor(198, 40, 40)
                tf.paragraphs[0].alignment = PP_ALIGN.RIGHT

                y += 1.1

        add_content_slide("⚠️ 개선 필요 영상", failure_content)

        # 슬라이드 6: 추천사항
        def recommendations_content(slide):
            recs = recommendations[:4]
            if not recs:
                tb = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
                tf = tb.text_frame
                tf.paragraphs[0].text = "추천사항 데이터가 없습니다."
                tf.paragraphs[0].font.size = Pt(18)
                return

            colors = {
                'critical': RGBColor(255, 235, 238),
                'high': RGBColor(255, 243, 224),
                'medium': RGBColor(227, 242, 253),
            }
            border_colors = {
                'critical': RGBColor(198, 40, 40),
                'high': RGBColor(230, 81, 0),
                'medium': RGBColor(21, 101, 192),
            }

            y = 1.5
            for rec in recs:
                priority = rec.get('priority', 'medium')
                category = rec.get('category', '')
                suggestions = rec.get('suggestions', [])[:2]

                box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12.3), Inches(1.3))
                box.fill.solid()
                box.fill.fore_color.rgb = colors.get(priority, colors['medium'])
                box.line.color.rgb = border_colors.get(priority, border_colors['medium'])

                cat_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.15), Inches(11.9), Inches(0.4))
                tf = cat_box.text_frame
                tf.paragraphs[0].text = category
                tf.paragraphs[0].font.size = Pt(16)
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.color.rgb = border_colors.get(priority, border_colors['medium'])

                sug_box = slide.shapes.add_textbox(Inches(0.7), Inches(y + 0.55), Inches(11.9), Inches(0.7))
                tf = sug_box.text_frame
                tf.word_wrap = True
                tf.paragraphs[0].text = " | ".join(suggestions)
                tf.paragraphs[0].font.size = Pt(13)
                tf.paragraphs[0].font.color.rgb = RGBColor(66, 66, 66)

                y += 1.4

        add_content_slide("💡 성장 전략 추천", recommendations_content)

        # PPT 바이트로 반환
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()

    def _generate_pdf_optimized_html(self) -> str:
        """PDF 최적화 HTML - 화면과 100% 동일한 내용 + 모든 영상 포함"""
        channel_summary = self.analysis.get('channel_summary', {})
        classification_stats = self.analysis.get('classification_stats', {})
        success = self.analysis.get('success_analysis', {})
        failure = self.analysis.get('failure_analysis', {})
        recommendations = self.analysis.get('recommendations', [])
        algorithm_insights = self.analysis.get('algorithm_insights', {})

        # 전체 영상 수집
        all_videos = []
        for cls in ['viral', 'hit', 'average', 'underperform']:
            videos = classification_stats.get(cls, {}).get('videos', [])
            for v in videos:
                v['classification'] = cls
                all_videos.append(v)
        all_videos.sort(key=lambda x: x.get('algorithm_score', 0), reverse=True)

        # 분류 데이터
        health = algorithm_insights.get('overall_health', '분석 중')
        avg_engagement = channel_summary.get('avg_engagement_rate', 0)
        avg_velocity = channel_summary.get('avg_view_velocity', 0)

        # 성공/실패 패턴 - 화면과 동일하게 5개만
        success_patterns = success.get('success_patterns', [])[:5]
        failure_patterns = failure.get('failure_patterns', [])[:5]

        # 분류별 카드 데이터
        viral_stats = classification_stats.get('viral', {})
        hit_stats = classification_stats.get('hit', {})
        avg_stats = classification_stats.get('average', {})
        under_stats = classification_stats.get('underperform', {})

        # 추천사항 HTML - 화면과 동일
        rec_html = ''
        for rec in recommendations:
            priority = rec.get('priority', 'medium')
            bg = '#ffebee' if priority in ['critical', 'high'] else '#fff3e0' if priority == 'medium' else '#e3f2fd'
            border = '#c62828' if priority in ['critical', 'high'] else '#e65100' if priority == 'medium' else '#1565c0'
            icon = '[긴급]' if priority in ['critical', 'high'] else '[중요]' if priority == 'medium' else '[참고]'
            current = rec.get('current', '')
            target = rec.get('target', '')
            meta = f'<p style="font-size:12px;color:#666;margin:0 0 8px 0;">현재: {current} -> 목표: {target}</p>' if current and target else ''
            rec_html += f'''
            <div style="background:{bg};border-left:4px solid {border};padding:15px;margin-bottom:12px;border-radius:0 8px 8px 0;page-break-inside:avoid;">
                <h4 style="margin:0 0 8px 0;color:{border};font-size:14px;">{icon} {rec.get('category', '')}</h4>
                {meta}
                <ul style="margin:0;padding-left:20px;font-size:13px;color:#333;">
                    {''.join(f"<li style='margin-bottom:5px;'>{s}</li>" for s in rec.get('suggestions', []))}
                </ul>
            </div>'''

        # 전체 영상 카드 HTML - 화면과 동일한 카드 형태
        class_colors = {'viral': '#d32f2f', 'hit': '#f57c00', 'average': '#1976d2', 'underperform': '#757575'}
        class_labels = {'viral': '바이럴', 'hit': '히트', 'average': '평균', 'underperform': '저조'}
        class_bg = {'viral': '#ffebee', 'hit': '#fff3e0', 'average': '#e3f2fd', 'underperform': '#f5f5f5'}

        videos_html = ''
        for v in all_videos:
            cls = v.get('classification', 'average')
            color = class_colors.get(cls, '#757575')
            label = class_labels.get(cls, '📊 평균')
            bg = class_bg.get(cls, '#f5f5f5')
            score = v.get('algorithm_score', 0)
            score_pct = min(100, score)
            videos_html += f'''
            <div style="display:inline-block;width:180px;margin:8px;vertical-align:top;background:#fff;border-radius:12px;box-shadow:0 2px 8px rgba(0,0,0,0.1);overflow:hidden;page-break-inside:avoid;">
                <div style="background:{bg};padding:8px;text-align:center;">
                    <span style="background:{color};color:#fff;padding:3px 10px;border-radius:10px;font-size:11px;font-weight:bold;">{label}</span>
                </div>
                <div style="padding:12px;">
                    <div style="font-size:11px;font-weight:600;height:32px;overflow:hidden;margin-bottom:8px;line-height:1.4;">{v.get('title', '제목 없음')[:40]}</div>
                    <div style="display:flex;justify-content:space-between;font-size:10px;color:#666;margin-bottom:4px;">
                        <span>조회수</span><span style="font-weight:600;color:#333;">{self._format_number(v.get('view_count', 0))}</span>
                    </div>
                    <div style="display:flex;justify-content:space-between;font-size:10px;color:#666;margin-bottom:4px;">
                        <span>참여율</span><span style="font-weight:600;color:#333;">{v.get('engagement_rate', 0):.2f}%</span>
                    </div>
                    <div style="display:flex;justify-content:space-between;font-size:10px;color:#666;margin-bottom:4px;">
                        <span>조회속도</span><span style="font-weight:600;color:#333;">{v.get('view_velocity', 0):.0f}/일</span>
                    </div>
                    <div style="display:flex;justify-content:space-between;font-size:10px;color:#666;margin-bottom:8px;">
                        <span>좋아요</span><span style="font-weight:600;color:#333;">{v.get('like_ratio', 0):.2f}%</span>
                    </div>
                    <div style="background:#e0e0e0;border-radius:4px;height:6px;overflow:hidden;">
                        <div style="background:{color};width:{score_pct}%;height:100%;"></div>
                    </div>
                    <div style="text-align:center;font-size:10px;color:{color};font-weight:bold;margin-top:4px;">점수: {score:.0f}</div>
                </div>
            </div>'''

        html = f'''<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8"/>
    <title>{channel_summary.get('channel_name', '')} - YouTube 분석 보고서</title>
    <style>
        @page {{ size: A4; margin: 15mm; }}
        * {{ box-sizing: border-box; }}
        body {{
            font-family: 'Malgun Gothic', 'Apple SD Gothic Neo', sans-serif;
            margin: 0;
            padding: 0;
            color: #333;
            background: #fff;
            font-size: 13px;
            line-height: 1.5;
        }}
        .section {{
            margin-bottom: 25px;
            page-break-inside: avoid;
        }}
        .section-title {{
            font-size: 16px;
            font-weight: 700;
            color: #212121;
            border-bottom: 3px solid #FF0000;
            padding-bottom: 8px;
            margin-bottom: 15px;
        }}
        .metrics-grid {{
            display: flex;
            gap: 10px;
            margin-bottom: 20px;
        }}
        .metric-card {{
            flex: 1;
            text-align: center;
            padding: 15px 10px;
            border-radius: 10px;
        }}
        .metric-card .icon {{ font-size: 20px; margin-bottom: 5px; }}
        .metric-card .value {{ font-size: 20px; font-weight: 700; }}
        .metric-card .label {{ font-size: 11px; color: #666; margin-top: 4px; }}
        .class-grid {{
            display: flex;
            gap: 12px;
            margin-bottom: 20px;
        }}
        .class-card {{
            flex: 1;
            text-align: center;
            padding: 20px 15px;
            border-radius: 12px;
        }}
        .class-card .count {{ font-size: 32px; font-weight: 700; }}
        .class-card .name {{ font-size: 13px; font-weight: 600; margin: 8px 0 4px; }}
        .class-card .avg {{ font-size: 11px; color: #666; }}
        .pattern-grid {{
            display: flex;
            gap: 15px;
            margin-bottom: 20px;
        }}
        .pattern-box {{
            flex: 1;
            padding: 20px;
            border-radius: 12px;
        }}
        .pattern-box h3 {{ margin: 0 0 12px 0; font-size: 14px; }}
        .pattern-box ul {{ margin: 0; padding-left: 20px; font-size: 12px; }}
        .pattern-box li {{ margin-bottom: 6px; }}
        .video-grid {{
            text-align: center;
        }}
    </style>
</head>
<body>

<!-- 알고리즘 건강도 -->
<div style="background: #CC0000; color: #fff; padding: 25px; border-radius: 16px; margin-bottom: 25px; text-align: center;">
    <div style="font-size: 13px; color: #ffcccc; margin-bottom: 8px;">YouTube 알고리즘 건강도</div>
    <div style="font-size: 36px; font-weight: 700; color: #ffffff; margin-bottom: 15px;">{health}</div>
    <div style="display: flex; justify-content: center; gap: 50px;">
        <div>
            <div style="font-size: 22px; font-weight: 700; color: #ffffff;">{avg_engagement:.2f}%</div>
            <div style="font-size: 11px; color: #ffcccc;">평균 참여율</div>
        </div>
        <div>
            <div style="font-size: 22px; font-weight: 700; color: #ffffff;">{self._format_number(channel_summary.get('avg_views_per_video', 0))}</div>
            <div style="font-size: 11px; color: #ffcccc;">평균 조회수</div>
        </div>
        <div>
            <div style="font-size: 22px; font-weight: 700; color: #ffffff;">{avg_velocity:.0f}/일</div>
            <div style="font-size: 11px; color: #ffcccc;">조회 속도</div>
        </div>
    </div>
</div>

<!-- 채널 정보 헤더 -->
<div style="background: #212121; color: #fff; padding: 20px; border-radius: 12px; margin-bottom: 25px;">
    <div style="background: #FF0000; display: inline-block; padding: 4px 12px; font-size: 10px; border-radius: 4px; margin-bottom: 10px; font-weight: 600;">YouTube Analytics Report v3.0</div>
    <h1 style="font-size: 22px; margin: 8px 0; font-weight: 700;">{channel_summary.get('channel_name', '채널')} 분석 보고서</h1>
    <div style="font-size: 12px; color: #bbb;">분석일: {datetime.now().strftime('%Y년 %m월 %d일')} | 분석 영상: {channel_summary.get('total_videos_analyzed', 0)}개 | 구독자: {channel_summary.get('subscriber_count', 0):,}명</div>
</div>

<!-- 핵심 지표 -->
<div class="section">
    <div class="section-title">[핵심 지표]</div>
    <div class="metrics-grid">
        <div class="metric-card" style="background:#fff3e0;">
            <div class="value" style="color:#e65100;">{self._format_number(channel_summary.get('subscriber_count', 0))}</div>
            <div class="label">구독자</div>
        </div>
        <div class="metric-card" style="background:#e8f5e9;">
            <div class="value" style="color:#2e7d32;">{channel_summary.get('total_videos_analyzed', 0)}</div>
            <div class="label">분석 영상</div>
        </div>
        <div class="metric-card" style="background:#e3f2fd;">
            <div class="value" style="color:#1565c0;">{self._format_number(channel_summary.get('avg_views_per_video', 0))}</div>
            <div class="label">평균 조회수</div>
        </div>
        <div class="metric-card" style="background:#f3e5f5;">
            <div class="value" style="color:#7b1fa2;">{avg_engagement:.2f}%</div>
            <div class="label">평균 참여율</div>
        </div>
        <div class="metric-card" style="background:#fce4ec;">
            <div class="value" style="color:#c62828;">{channel_summary.get('avg_like_ratio', 0):.2f}%</div>
            <div class="label">좋아요 비율</div>
        </div>
    </div>
</div>

<!-- 영상 성과 분류 -->
<div class="section">
    <div class="section-title">[영상 성과 분류]</div>
    <div class="class-grid">
        <div class="class-card" style="background:#ffebee;border:2px solid #d32f2f;">
            <div class="count" style="color:#d32f2f;">{viral_stats.get('count', 0)}</div>
            <div class="name" style="color:#d32f2f;">바이럴</div>
            <div class="avg">평균 {self._format_number(viral_stats.get('avg_views', 0))} 조회</div>
        </div>
        <div class="class-card" style="background:#fff3e0;border:2px solid #f57c00;">
            <div class="count" style="color:#f57c00;">{hit_stats.get('count', 0)}</div>
            <div class="name" style="color:#f57c00;">히트</div>
            <div class="avg">평균 {self._format_number(hit_stats.get('avg_views', 0))} 조회</div>
        </div>
        <div class="class-card" style="background:#e3f2fd;border:2px solid #1976d2;">
            <div class="count" style="color:#1976d2;">{avg_stats.get('count', 0)}</div>
            <div class="name" style="color:#1976d2;">평균</div>
            <div class="avg">평균 {self._format_number(avg_stats.get('avg_views', 0))} 조회</div>
        </div>
        <div class="class-card" style="background:#f5f5f5;border:2px solid #757575;">
            <div class="count" style="color:#757575;">{under_stats.get('count', 0)}</div>
            <div class="name" style="color:#757575;">저조</div>
            <div class="avg">평균 {self._format_number(under_stats.get('avg_views', 0))} 조회</div>
        </div>
    </div>
</div>

<!-- 성공/실패 패턴 -->
<div class="section">
    <div class="section-title">[성공/개선 패턴 분석]</div>
    <div class="pattern-grid">
        <div class="pattern-box" style="background:#e8f5e9; border: 2px solid #4caf50;">
            <h3 style="color:#2e7d32;">[성공] 성공 패턴</h3>
            <ul style="color:#333;">
                {''.join(f"<li>{p}</li>" for p in success_patterns) if success_patterns else '<li>분석 중...</li>'}
            </ul>
        </div>
        <div class="pattern-box" style="background:#ffebee; border: 2px solid #f44336;">
            <h3 style="color:#c62828;">[주의] 개선 필요</h3>
            <ul style="color:#333;">
                {''.join(f"<li>{p}</li>" for p in failure_patterns) if failure_patterns else '<li>개선 사항 없음</li>'}
            </ul>
        </div>
    </div>
</div>

<!-- 추천 전략 -->
<div class="section">
    <div class="section-title">[데이터 기반 추천 전략]</div>
    {rec_html if rec_html else '<p style="color:#666;">추천 데이터 없음</p>'}
</div>

<!-- 전체 영상 목록 -->
<div class="section" style="page-break-before: always;">
    <div class="section-title">[전체 영상 성과] ({len(all_videos)}개)</div>
    <div style="background:#f5f5f5;padding:10px 15px;margin-bottom:15px;border-radius:8px;font-size:11px;">
        <strong>분류 기준:</strong>
        <span style="background:#d32f2f;color:#fff;padding:2px 10px;border-radius:10px;margin-left:10px;">바이럴</span>
        <span style="background:#f57c00;color:#fff;padding:2px 10px;border-radius:10px;margin-left:8px;">히트</span>
        <span style="background:#1976d2;color:#fff;padding:2px 10px;border-radius:10px;margin-left:8px;">평균</span>
        <span style="background:#757575;color:#fff;padding:2px 10px;border-radius:10px;margin-left:8px;">저조</span>
    </div>
    <div class="video-grid">
        {videos_html}
    </div>
</div>

<!-- 푸터 -->
<div style="text-align:center;padding:20px;border-top:2px solid #e0e0e0;margin-top:30px;color:#999;font-size:11px;">
    YouTube Analytics Report v3.0 | {channel_summary.get('channel_name', '')} | 생성일: {datetime.now().strftime('%Y-%m-%d %H:%M')}
</div>

</body>
</html>'''
        return html


class CompetitorReportGenerator:
    """경쟁사 비교 보고서 생성기 - 상세 버전"""

    THEME = ReportGenerator.THEME

    def __init__(self, comparison_data: dict):
        self.comparison = comparison_data

    def generate_html_report(self) -> str:
        """경쟁사 비교 HTML 보고서 - 완전 상세 버전"""
        main = self.comparison.get('main_channel', {})
        competitors = self.comparison.get('competitors', [])
        position = self.comparison.get('market_position', {})
        strengths = self.comparison.get('strengths_weaknesses', {})
        recommendations = self.comparison.get('recommendations', [])
        rankings = self.comparison.get('ranking_analysis', {})
        strategy = self.comparison.get('content_strategy_comparison', {})
        gaps = self.comparison.get('performance_gap_analysis', {})
        insights = self.comparison.get('competitive_insights', [])
        metrics = self.comparison.get('metrics_comparison', {})

        # 순위 테이블 HTML
        ranking_rows = ''
        for key, data in rankings.items():
            if isinstance(data, dict):
                rank = data.get('rank', '-')
                total = data.get('total', '-')
                name = data.get('name', key)
                value = data.get('value', 0)
                top_channel = data.get('top_channel', '-')
                gap = data.get('gap_to_top', 0)
                is_top = data.get('is_top', False)

                rank_color = '#d32f2f' if rank == 1 else '#f57c00' if rank == 2 else '#1976d2' if rank <= 3 else '#757575'
                rank_icon = '🥇' if rank == 1 else '🥈' if rank == 2 else '🥉' if rank == 3 else f'{rank}위'

                if isinstance(value, float):
                    value_str = f'{value:.1f}%' if 'rate' in key or 'engagement' in key else f'{value:.0f}'
                else:
                    value_str = self._format_number(value)

                ranking_rows += f'''
                <tr>
                    <td style="font-weight:600;">{name}</td>
                    <td style="text-align:center;color:{rank_color};font-weight:700;font-size:18px;">{rank_icon}</td>
                    <td style="text-align:center;">{total}개 중</td>
                    <td style="text-align:right;font-weight:600;">{value_str}</td>
                    <td style="text-align:center;">{top_channel if not is_top else '본인'}</td>
                    <td style="text-align:right;color:{'#2e7d32' if gap <= 0 else '#c62828'};">{'+' if gap < 0 else ''}{self._format_gap(gap, key) if gap != 0 else '-'}</td>
                </tr>'''

        # 인사이트 HTML
        insights_html = ''
        for ins in insights:
            bg = '#e8f5e9' if ins.get('type') == 'positive' else '#ffebee'
            icon = '✅' if ins.get('type') == 'positive' else '⚠️'
            color = '#2e7d32' if ins.get('type') == 'positive' else '#c62828'
            insights_html += f'''
            <div style="background:{bg};padding:16px;border-radius:10px;margin-bottom:12px;">
                <div style="font-weight:700;color:{color};margin-bottom:6px;">{icon} {ins.get('title', '')}</div>
                <div style="font-size:14px;color:#555;">{ins.get('detail', '')}</div>
            </div>'''

        # 전략 비교 HTML
        strategy_html = ''
        for comp_strat in strategy.get('competitor_strategies', []):
            strategy_html += f'''
            <tr>
                <td>{comp_strat.get('channel_name', '')}</td>
                <td style="text-align:center;font-weight:600;">{comp_strat.get('strategy', '')}</td>
                <td style="text-align:right;">{comp_strat.get('viral_rate', 0):.1f}%</td>
                <td style="text-align:right;">{comp_strat.get('engagement', 0):.2f}%</td>
            </tr>'''

        # 추천 전략 HTML
        rec_html = ''
        priority_colors = {'critical': '#c62828', 'high': '#e65100', 'medium': '#1976d2', 'low': '#757575'}
        priority_bg = {'critical': '#ffebee', 'high': '#fff3e0', 'medium': '#e3f2fd', 'low': '#f5f5f5'}
        priority_icons = {'critical': '🔴', 'high': '🟠', 'medium': '🔵', 'low': '⚪'}

        for r in recommendations:
            p = r.get('priority', 'medium')
            rec_html += f'''
            <div style="background:{priority_bg.get(p, '#f5f5f5')};border-left:4px solid {priority_colors.get(p, '#757575')};padding:20px;border-radius:0 12px 12px 0;margin-bottom:16px;">
                <h4 style="color:{priority_colors.get(p, '#757575')};margin-bottom:12px;font-size:16px;">{priority_icons.get(p, '🔵')} {r.get('category', '')}</h4>
                <ul style="padding-left:20px;margin:0;">
                    {''.join(f'<li style="margin-bottom:8px;font-size:14px;">{s}</li>' for s in r.get('suggestions', []))}
                </ul>
            </div>'''

        # 성과 격차 분석 HTML
        vs_avg = gaps.get('vs_average', {})
        vs_best = gaps.get('vs_best', {})
        advantages = gaps.get('competitive_advantages', [])

        html = f'''<!DOCTYPE html>
<html lang="ko">
<head>
    <meta charset="UTF-8">
    <title>경쟁사 비교 분석 보고서 - {main.get('channel_name', '')}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: 'Malgun Gothic', 'Apple SD Gothic Neo', sans-serif; background: #f5f5f5; color: #212121; line-height: 1.6; font-size: 14px; }}
        .container {{ max-width: 1200px; margin: 0 auto; padding: 30px 20px; }}
        .cover {{ background: linear-gradient(135deg, #FF0000, #CC0000); color: white; padding: 40px; border-radius: 16px; margin-bottom: 30px; }}
        .cover h1 {{ font-size: 28px; margin-bottom: 8px; }}
        .cover p {{ opacity: 0.9; font-size: 16px; }}
        .section {{ background: white; border-radius: 12px; padding: 25px; margin-bottom: 20px; box-shadow: 0 2px 8px rgba(0,0,0,0.06); }}
        .section-title {{ font-size: 18px; font-weight: 700; color: #212121; margin-bottom: 20px; padding-bottom: 10px; border-bottom: 3px solid #FF0000; display: flex; align-items: center; gap: 10px; }}
        .grid-2 {{ display: grid; grid-template-columns: 1fr 1fr; gap: 20px; }}
        .grid-3 {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 15px; }}
        .grid-4 {{ display: grid; grid-template-columns: repeat(4, 1fr); gap: 15px; }}
        .card {{ background: #f9f9f9; padding: 20px; border-radius: 10px; text-align: center; }}
        .card .value {{ font-size: 28px; font-weight: 700; color: #FF0000; }}
        .card .label {{ font-size: 12px; color: #666; margin-top: 5px; }}
        .card .sub {{ font-size: 11px; color: #999; margin-top: 3px; }}
        table {{ width: 100%; border-collapse: collapse; }}
        th {{ background: #212121; color: white; padding: 12px; text-align: left; font-size: 13px; }}
        td {{ padding: 12px; border-bottom: 1px solid #e0e0e0; font-size: 13px; }}
        tr:hover {{ background: #fafafa; }}
        .highlight {{ background: #fff3e0 !important; }}
        .badge {{ display: inline-block; padding: 4px 12px; border-radius: 20px; font-size: 12px; font-weight: 600; }}
        .badge-red {{ background: #ffebee; color: #c62828; }}
        .badge-green {{ background: #e8f5e9; color: #2e7d32; }}
        .badge-blue {{ background: #e3f2fd; color: #1565c0; }}
        .progress-bar {{ background: #e0e0e0; border-radius: 10px; height: 10px; overflow: hidden; }}
        .progress-fill {{ height: 100%; border-radius: 10px; }}
    </style>
</head>
<body>
<div class="container">
    <!-- 표지 -->
    <div class="cover">
        <h1>📊 경쟁사 비교 분석 보고서</h1>
        <p>{main.get('channel_name', '내 채널')} vs 경쟁 채널 {len(competitors)}개 심층 분석</p>
        <div style="margin-top:20px;display:flex;gap:30px;font-size:14px;">
            <div>분석 영상: {main.get('total_videos', 0)}개</div>
            <div>경쟁사: {len(competitors)}개 채널</div>
            <div>생성일: {datetime.now().strftime('%Y-%m-%d')}</div>
        </div>
    </div>

    <!-- 종합 점수 -->
    <div class="section">
        <h2 class="section-title">🏆 시장 포지션 & 종합 점수</h2>
        <div class="grid-2">
            <div style="text-align:center;padding:30px;background:linear-gradient(135deg,#fff3e0,#ffe0b2);border-radius:12px;">
                <div style="font-size:64px;">{'🥇' if position.get('position') == 'dominant' else '🥈' if position.get('position') == 'leader' else '🥉' if position.get('position') == 'challenger' else '🏅'}</div>
                <div style="font-size:28px;font-weight:700;color:#e65100;margin:10px 0;">{position.get('position_korean', '분석 중')}</div>
                <div style="font-size:14px;color:#666;">{position.get('interpretation', '')}</div>
            </div>
            <div>
                <div style="font-size:14px;font-weight:600;margin-bottom:15px;">경쟁력 점수: <span style="font-size:24px;color:#FF0000;">{position.get('overall_score', 0):.0f}</span>/100</div>
                <div style="margin-bottom:12px;">
                    <div style="display:flex;justify-content:space-between;font-size:12px;margin-bottom:4px;"><span>구독자</span><span>{position.get('subscriber_vs_avg', 0):.0f}%</span></div>
                    <div class="progress-bar"><div class="progress-fill" style="width:{min(100, position.get('subscriber_vs_avg', 0))}%;background:#FF0000;"></div></div>
                </div>
                <div style="margin-bottom:12px;">
                    <div style="display:flex;justify-content:space-between;font-size:12px;margin-bottom:4px;"><span>조회수</span><span>{position.get('views_vs_avg', 0):.0f}%</span></div>
                    <div class="progress-bar"><div class="progress-fill" style="width:{min(100, position.get('views_vs_avg', 0))}%;background:#e65100;"></div></div>
                </div>
                <div style="margin-bottom:12px;">
                    <div style="display:flex;justify-content:space-between;font-size:12px;margin-bottom:4px;"><span>참여율</span><span>{position.get('engagement_vs_avg', 0):.0f}%</span></div>
                    <div class="progress-bar"><div class="progress-fill" style="width:{min(100, position.get('engagement_vs_avg', 0))}%;background:#1976d2;"></div></div>
                </div>
                <div>
                    <div style="display:flex;justify-content:space-between;font-size:12px;margin-bottom:4px;"><span>조회속도</span><span>{position.get('velocity_vs_avg', 0):.0f}%</span></div>
                    <div class="progress-bar"><div class="progress-fill" style="width:{min(100, position.get('velocity_vs_avg', 0))}%;background:#2e7d32;"></div></div>
                </div>
            </div>
        </div>
    </div>

    <!-- 핵심 인사이트 -->
    {f'''<div class="section">
        <h2 class="section-title">💡 핵심 경쟁 인사이트</h2>
        {insights_html if insights_html else '<p style="color:#666;">인사이트 분석 중...</p>'}
    </div>''' if insights else ''}

    <!-- 지표별 순위 -->
    <div class="section">
        <h2 class="section-title">📊 지표별 순위 분석</h2>
        <table>
            <thead>
                <tr>
                    <th>지표</th>
                    <th style="text-align:center;">순위</th>
                    <th style="text-align:center;">전체</th>
                    <th style="text-align:right;">내 수치</th>
                    <th style="text-align:center;">1위 채널</th>
                    <th style="text-align:right;">1위와 격차</th>
                </tr>
            </thead>
            <tbody>
                {ranking_rows if ranking_rows else '<tr><td colspan="6" style="text-align:center;color:#666;">순위 데이터 없음</td></tr>'}
            </tbody>
        </table>
    </div>

    <!-- 성과 격차 분석 -->
    <div class="section">
        <h2 class="section-title">📈 경쟁사 대비 성과 격차</h2>
        <div class="grid-2">
            <div>
                <h4 style="margin-bottom:15px;color:#1976d2;">vs 경쟁사 평균</h4>
                <div class="grid-2" style="gap:10px;">
                    <div class="card">
                        <div class="value" style="color:{'#2e7d32' if vs_avg.get('subscriber_gap_pct', 0) >= 0 else '#c62828'};">{'+' if vs_avg.get('subscriber_gap_pct', 0) >= 0 else ''}{vs_avg.get('subscriber_gap_pct', 0):.0f}%</div>
                        <div class="label">구독자</div>
                    </div>
                    <div class="card">
                        <div class="value" style="color:{'#2e7d32' if vs_avg.get('views_gap_pct', 0) >= 0 else '#c62828'};">{'+' if vs_avg.get('views_gap_pct', 0) >= 0 else ''}{vs_avg.get('views_gap_pct', 0):.0f}%</div>
                        <div class="label">조회수</div>
                    </div>
                    <div class="card">
                        <div class="value" style="color:{'#2e7d32' if vs_avg.get('engagement_gap', 0) >= 0 else '#c62828'};">{'+' if vs_avg.get('engagement_gap', 0) >= 0 else ''}{vs_avg.get('engagement_gap', 0):.2f}%p</div>
                        <div class="label">참여율</div>
                    </div>
                    <div class="card">
                        <div class="value" style="color:{'#2e7d32' if vs_avg.get('viral_rate_gap', 0) >= 0 else '#c62828'};">{'+' if vs_avg.get('viral_rate_gap', 0) >= 0 else ''}{vs_avg.get('viral_rate_gap', 0):.1f}%p</div>
                        <div class="label">바이럴율</div>
                    </div>
                </div>
            </div>
            <div>
                <h4 style="margin-bottom:15px;color:#e65100;">vs 최고 성과 채널</h4>
                <div style="font-size:13px;">
                    <div style="padding:10px;background:#f5f5f5;border-radius:8px;margin-bottom:8px;">
                        <span style="color:#666;">구독자 1위:</span> <strong>{vs_best.get('best_subscriber_channel', '-')}</strong>
                        <span style="float:right;color:{'#2e7d32' if vs_best.get('subscriber_gap_to_best', 0) >= 0 else '#c62828'};">{self._format_number(vs_best.get('subscriber_gap_to_best', 0))} 차이</span>
                    </div>
                    <div style="padding:10px;background:#f5f5f5;border-radius:8px;margin-bottom:8px;">
                        <span style="color:#666;">조회수 1위:</span> <strong>{vs_best.get('best_views_channel', '-')}</strong>
                        <span style="float:right;color:{'#2e7d32' if vs_best.get('views_gap_to_best', 0) >= 0 else '#c62828'};">{self._format_number(vs_best.get('views_gap_to_best', 0))} 차이</span>
                    </div>
                    <div style="padding:10px;background:#f5f5f5;border-radius:8px;margin-bottom:8px;">
                        <span style="color:#666;">참여율 1위:</span> <strong>{vs_best.get('best_engagement_channel', '-')}</strong>
                        <span style="float:right;color:{'#2e7d32' if vs_best.get('engagement_gap_to_best', 0) >= 0 else '#c62828'};">{vs_best.get('engagement_gap_to_best', 0):+.2f}%p</span>
                    </div>
                    <div style="padding:10px;background:#f5f5f5;border-radius:8px;">
                        <span style="color:#666;">바이럴 1위:</span> <strong>{vs_best.get('best_viral_channel', '-')}</strong>
                        <span style="float:right;color:{'#2e7d32' if vs_best.get('viral_gap_to_best', 0) >= 0 else '#c62828'};">{vs_best.get('viral_gap_to_best', 0):+.1f}%p</span>
                    </div>
                </div>
            </div>
        </div>
        {f'''<div style="margin-top:20px;">
            <h4 style="margin-bottom:10px;color:#2e7d32;">✅ 경쟁 우위 항목</h4>
            <div style="display:flex;flex-wrap:wrap;gap:8px;">
                {''.join(f'<span class="badge badge-green">{a}</span>' for a in advantages)}
            </div>
        </div>''' if advantages else ''}
    </div>

    <!-- 콘텐츠 전략 비교 -->
    <div class="section">
        <h2 class="section-title">🎯 콘텐츠 전략 비교</h2>
        <div style="text-align:center;padding:20px;background:#f5f5f5;border-radius:12px;margin-bottom:20px;">
            <div style="font-size:14px;color:#666;">내 채널 전략 유형</div>
            <div style="font-size:24px;font-weight:700;color:#FF0000;margin-top:8px;">{strategy.get('main_strategy', '분석 중')}</div>
        </div>
        <table>
            <thead>
                <tr>
                    <th>채널</th>
                    <th style="text-align:center;">전략 유형</th>
                    <th style="text-align:right;">바이럴율</th>
                    <th style="text-align:right;">참여율</th>
                </tr>
            </thead>
            <tbody>
                <tr class="highlight">
                    <td><strong>{main.get('channel_name', '내 채널')}</strong></td>
                    <td style="text-align:center;font-weight:600;">{strategy.get('main_strategy', '-')}</td>
                    <td style="text-align:right;">{main.get('viral_rate', 0):.1f}%</td>
                    <td style="text-align:right;">{main.get('avg_engagement', 0):.2f}%</td>
                </tr>
                {strategy_html}
            </tbody>
        </table>
    </div>

    <!-- 강점/약점 분석 -->
    <div class="section">
        <h2 class="section-title">💪 강점 & ⚠️ 약점 분석</h2>
        <div class="grid-2">
            <div style="background:#e8f5e9;padding:25px;border-radius:12px;">
                <h3 style="color:#2e7d32;margin-bottom:15px;font-size:16px;">💪 강점 ({strengths.get('strength_count', 0)}개)</h3>
                <ul style="padding-left:20px;margin:0;">
                    {''.join(f'<li style="margin-bottom:10px;font-size:14px;">{s}</li>' for s in strengths.get('strengths', ['분석 중']))}
                </ul>
            </div>
            <div style="background:#ffebee;padding:25px;border-radius:12px;">
                <h3 style="color:#c62828;margin-bottom:15px;font-size:16px;">⚠️ 개선 필요 ({strengths.get('weakness_count', 0)}개)</h3>
                <ul style="padding-left:20px;margin:0;">
                    {''.join(f'<li style="margin-bottom:10px;font-size:14px;">{w}</li>' for w in strengths.get('weaknesses', ['분석 중']))}
                </ul>
            </div>
        </div>
    </div>

    <!-- 채널별 상세 비교 -->
    <div class="section">
        <h2 class="section-title">📋 채널별 상세 비교표</h2>
        <table>
            <thead>
                <tr>
                    <th>채널</th>
                    <th style="text-align:right;">구독자</th>
                    <th style="text-align:right;">평균 조회수</th>
                    <th style="text-align:right;">조회 속도</th>
                    <th style="text-align:right;">참여율</th>
                    <th style="text-align:right;">바이럴</th>
                    <th style="text-align:right;">히트</th>
                    <th style="text-align:right;">성공률</th>
                </tr>
            </thead>
            <tbody>
                <tr class="highlight">
                    <td><strong>⭐ {main.get('channel_name', '내 채널')}</strong></td>
                    <td style="text-align:right;font-weight:600;">{self._format_number(main.get('subscriber_count', 0))}</td>
                    <td style="text-align:right;font-weight:600;">{self._format_number(main.get('avg_views', 0))}</td>
                    <td style="text-align:right;">{main.get('avg_velocity', 0):.0f}/일</td>
                    <td style="text-align:right;">{main.get('avg_engagement', 0):.2f}%</td>
                    <td style="text-align:right;color:#c62828;font-weight:600;">{main.get('viral_count', 0)}개</td>
                    <td style="text-align:right;color:#e65100;">{main.get('hit_count', 0)}개</td>
                    <td style="text-align:right;font-weight:600;">{main.get('success_rate', 0):.1f}%</td>
                </tr>
                {''.join(f"""
                <tr>
                    <td>{c.get('channel_name', '')}</td>
                    <td style="text-align:right;">{self._format_number(c.get('subscriber_count', 0))}</td>
                    <td style="text-align:right;">{self._format_number(c.get('avg_views', 0))}</td>
                    <td style="text-align:right;">{c.get('avg_velocity', 0):.0f}/일</td>
                    <td style="text-align:right;">{c.get('avg_engagement', 0):.2f}%</td>
                    <td style="text-align:right;">{c.get('viral_count', 0)}개</td>
                    <td style="text-align:right;">{c.get('hit_count', 0)}개</td>
                    <td style="text-align:right;">{c.get('success_rate', 0):.1f}%</td>
                </tr>
                """ for c in competitors)}
            </tbody>
        </table>
    </div>

    <!-- 경쟁 전략 추천 -->
    <div class="section">
        <h2 class="section-title">🚀 경쟁 전략 추천</h2>
        {rec_html if rec_html else '<p style="color:#666;">추천 전략 생성 중...</p>'}
    </div>

    <!-- 푸터 -->
    <div style="text-align:center;padding:20px;color:#999;font-size:12px;border-top:1px solid #e0e0e0;margin-top:20px;">
        YouTube 경쟁사 비교 분석 보고서 | {main.get('channel_name', '')} | {datetime.now().strftime('%Y-%m-%d %H:%M')}
    </div>
</div>
</body>
</html>'''
        return html

    def _format_number(self, num) -> str:
        if num is None:
            return '0'
        num = float(num)
        if num >= 1000000:
            return f'{num / 1000000:.1f}M'
        elif num >= 1000:
            return f'{num / 1000:.1f}K'
        return str(int(num))

    def _format_gap(self, gap, metric):
        if metric in ['engagement', 'avg_engagement', 'viral_rate', 'success_rate']:
            return f'{abs(gap):.1f}%p'
        elif abs(gap) >= 1000000:
            return f'{gap/1000000:+.1f}M'
        elif abs(gap) >= 1000:
            return f'{gap/1000:+.1f}K'
        return f'{int(gap):+d}'

    def generate_pdf_report(self) -> bytes:
        """경쟁사 비교 PDF 보고서 생성"""
        try:
            import pdfkit
            import shutil

            # wkhtmltopdf 경로 찾기
            wkhtmltopdf_path = shutil.which('wkhtmltopdf')
            if not wkhtmltopdf_path:
                for path in [
                    r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe',
                    r'C:\Program Files (x86)\wkhtmltopdf\bin\wkhtmltopdf.exe',
                    r'C:\ProgramData\chocolatey\bin\wkhtmltopdf.exe',
                ]:
                    if os.path.exists(path):
                        wkhtmltopdf_path = path
                        break

            config = pdfkit.configuration(wkhtmltopdf=wkhtmltopdf_path) if wkhtmltopdf_path else None

            options = {
                'encoding': 'UTF-8',
                'page-size': 'A4',
                'margin-top': '15mm',
                'margin-right': '15mm',
                'margin-bottom': '15mm',
                'margin-left': '15mm',
                'enable-local-file-access': None,
                'no-stop-slow-scripts': None,
            }

            html_content = self.generate_html_report()
            pdf_bytes = pdfkit.from_string(html_content, False, options=options, configuration=config)
            return pdf_bytes

        except Exception as e:
            raise ImportError(f"PDF 생성 실패: {str(e)}")

    def generate_json_report(self) -> str:
        """경쟁사 비교 JSON 보고서"""
        import json
        return json.dumps(self.comparison, ensure_ascii=False, indent=2)
